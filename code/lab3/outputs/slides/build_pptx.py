"""Render a deck JSON spec into an editable .pptx.

  python3 build_pptx.py deck_final.json lab3_journey.pptx

Each slide: assertion title (English), bullets (English), one figure, and the
Traditional-Chinese speaker note placed in PowerPoint's notes pane.
Figures are fit into their box preserving aspect ratio.
"""
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
LAB3 = HERE.parents[1]            # code/lab3
FIG_ROOT = LAB3 / 'outputs' / 'figures'

# 16:9 canvas
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

NAVY = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x1B, 0x6E, 0xC2)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)


def fit_box(img_path, bx, by, bw, bh):
    """Return (left, top, width, height) EMU fitting the image in the box, centered."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(bw / iw, bh / ih)
    w = int(iw * scale); h = int(ih * scale)
    left = int(bx + (bw - w) / 2)
    top = int(by + (bh - h) / 2)
    return left, top, w, h


def add_title(slide, text, top=Inches(0.35), height=Inches(0.95)):
    tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(27); r.font.bold = True; r.font.color.rgb = NAVY
    return tb


def add_bullets(slide, bullets, left, top, width, height, size=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        dot = p.add_run(); dot.text = '•  '
        dot.font.size = Pt(size); dot.font.color.rgb = ACCENT; dot.font.bold = True
        r = p.add_run(); r.text = b
        r.font.size = Pt(size); r.font.color.rgb = NAVY
    return tb


def resolve_fig(figure):
    if not figure:
        return None
    p = FIG_ROOT / figure
    return p if p.exists() else None


def set_notes(slide, notes):
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def build(spec_path, out_path):
    spec = json.load(open(spec_path, encoding='utf-8'))
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    slides = spec['slides']
    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        is_title = (i == 0) or s.get('id') == 'title'
        fig = resolve_fig(s.get('figure'))
        bullets = s.get('bullets') or []

        if is_title:
            # centered title + subtitle
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(1.6))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = s['title']
            r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY
            sub = s.get('subtitle')
            if sub:
                tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.33), Inches(1.0))
                tf2 = tb2.text_frame; tf2.word_wrap = True
                p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run(); r2.text = sub
                r2.font.size = Pt(20); r2.font.color.rgb = GREY
            set_notes(slide, s.get('notes', ''))
            continue

        add_title(slide, s['title'])

        if fig is None:
            # text-only: centered bullets, larger
            add_bullets(slide, bullets, Inches(1.5), Inches(1.6),
                        Inches(10.3), Inches(5.4), size=22)
        else:
            with Image.open(fig) as im:
                iw, ih = im.size
            aspect = iw / ih
            if aspect >= 1.5 and bullets:
                # wide figure: bullets band under title, image fills the rest
                add_bullets(slide, bullets, Inches(0.6), Inches(1.25),
                            Inches(12.1), Inches(1.55), size=16)
                l, t, w, h = fit_box(fig, Inches(0.5), Inches(2.95),
                                     Inches(12.33), Inches(4.3))
                slide.shapes.add_picture(str(fig), l, t, w, h)
            elif not bullets:
                # figure only: large centered
                l, t, w, h = fit_box(fig, Inches(0.5), Inches(1.3),
                                     Inches(12.33), Inches(5.9))
                slide.shapes.add_picture(str(fig), l, t, w, h)
            else:
                # tall/square: bullets left, image right
                add_bullets(slide, bullets, Inches(0.6), Inches(1.4),
                            Inches(4.5), Inches(5.6), size=18)
                l, t, w, h = fit_box(fig, Inches(5.3), Inches(1.3),
                                     Inches(7.6), Inches(5.9))
                slide.shapes.add_picture(str(fig), l, t, w, h)

        set_notes(slide, s.get('notes', ''))

    prs.save(out_path)
    print(f'  ✓ {out_path}  ({len(slides)} slides)')


if __name__ == '__main__':
    spec = sys.argv[1] if len(sys.argv) > 1 else str(HERE / 'deck_final.json')
    out = sys.argv[2] if len(sys.argv) > 2 else str(HERE / 'lab3_journey.pptx')
    build(spec, out)
