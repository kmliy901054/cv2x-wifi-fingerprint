"""Render a deck JSON spec into an editable .pptx.

  python3 build_pptx.py deck_final.json lab3_journey.pptx

Each slide: assertion title (English), bullets (English), one figure, and the
Traditional-Chinese speaker note placed in PowerPoint's notes pane.
Figures are fit into their box preserving aspect ratio.
"""
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
LAB3 = HERE.parents[1]            # lab3
FIG_ROOT = LAB3 / 'outputs' / 'figures'

# 16:9 canvas
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

NAVY = RGBColor(0x1F, 0x2D, 0x3D)
ACCENT = RGBColor(0x1B, 0x6E, 0xC2)
ACCENT_LT = RGBColor(0x6F, 0xB7, 0xF0)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def fit_box(img_path, bx, by, bw, bh):
    """Return (left, top, width, height) EMU fitting the image in the box, centered."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(bw / iw, bh / ih)
    w = int(iw * scale); h = int(ih * scale)
    left = int(bx + (bw - w) / 2)
    top = int(by + (bh - h) / 2)
    return left, top, w, h


def add_title(slide, text, top=Inches(0.35), height=Inches(0.95)):
    tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.33), height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(27); r.font.bold = True; r.font.color.rgb = NAVY
    return tb


def add_bullets(slide, bullets, left, top, width, height, size=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        dot = p.add_run(); dot.text = '•  '
        dot.font.size = Pt(size); dot.font.color.rgb = ACCENT; dot.font.bold = True
        r = p.add_run(); r.text = b
        r.font.size = Pt(size); r.font.color.rgb = NAVY
    return tb


def resolve_fig(figure):
    if not figure:
        return None
    p = FIG_ROOT / figure
    return p if p.exists() else None


def set_notes(slide, notes):
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def add_play_button(slide, l, t, w, h):
    """Overlay a play-button + caption on an image box to mark a video placeholder."""
    from pptx.enum.shapes import MSO_SHAPE
    cx, cy = l + w // 2, t + h // 2
    d = min(w, h) // 5
    # dark translucent circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d)
    circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x11, 0x11, 0x11)
    circ.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); circ.line.width = Pt(2)
    try:
        circ.fill.transparency = 0.25
    except Exception:
        pass
    # white play triangle
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  cx - d // 8, cy - d // 4, d // 2, d // 2)
    tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tri.line.fill.background()
    # caption ribbon under the image
    cap = slide.shapes.add_textbox(l, t + h - Inches(0.45), w, Inches(0.45))
    cf = cap.text_frame; cf.word_wrap = True
    p = cf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '▶  LIVE DEMO VIDEO — replace this slide image with the recording'
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT


def add_bg(slide, color):
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_rule(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def build_title_slide(slide, s):
    add_bg(slide, NAVY)
    add_rule(slide, Inches(5.17), Inches(2.35), Inches(3.0), Pt(5), ACCENT_LT)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.73), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s['title']
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    sub = s.get('subtitle')
    if sub:
        tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.33), Inches(1.2))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(20); r2.font.color.rgb = ACCENT_LT
    auth = s.get('author')
    if auth:
        tb3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.33), Inches(0.6))
        p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = auth
        r3.font.size = Pt(14); r3.font.color.rgb = RGBColor(0xAA, 0xB6, 0xC4)


def build_divider_slide(slide, s, idx, total):
    add_bg(slide, NAVY)
    # kicker: SECTION n / total
    kb = slide.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.6))
    kp = kb.text_frame.paragraphs[0]; kp.alignment = PP_ALIGN.CENTER
    kr = kp.add_run(); kr.text = f'SECTION {idx} / {total}'
    kr.font.size = Pt(18); kr.font.bold = True; kr.font.color.rgb = ACCENT_LT
    # big section title (centered)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.05), Inches(11.73), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = s['title']
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    add_rule(slide, Inches(5.67), Inches(4.55), Inches(2.0), Pt(4), ACCENT_LT)
    sub = s.get('subtitle')
    if sub:
        sb = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.33), Inches(1.0))
        sf = sb.text_frame; sf.word_wrap = True
        sp = sf.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run(); sr.text = sub
        sr.font.size = Pt(18); sr.font.color.rgb = RGBColor(0xC8, 0xD2, 0xDE)


def build_textcard_slide(slide, s, bullets):
    """Text-only content slide: title + a filled card holding large bullets."""
    add_title(slide, s['title'])
    sub = s.get('subtitle')
    top = Inches(1.45)
    if sub:
        sb = slide.shapes.add_textbox(Inches(0.6), Inches(1.28), Inches(12.1), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]
        sr = sp.add_run(); sr.text = sub
        sr.font.size = Pt(16); sr.font.italic = True; sr.font.color.rgb = ACCENT
        top = Inches(1.85)
    # card
    card = add_rule(slide, Inches(0.7), top, Inches(11.93), Inches(7.15) - top, LIGHT)
    add_rule(slide, Inches(0.7), top, Pt(7), Inches(7.15) - top, ACCENT)  # left accent bar
    # bullets inside card, generously sized, top-anchored
    tb = slide.shapes.add_textbox(Inches(1.15), top + Inches(0.35),
                                   Inches(11.1), Inches(7.15) - top - Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    n = len(bullets)
    size = 26 if n <= 4 else (22 if n <= 6 else 18)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(16)
        dot = p.add_run(); dot.text = '▪  '
        dot.font.size = Pt(size); dot.font.color.rgb = ACCENT; dot.font.bold = True
        r = p.add_run(); r.text = b
        r.font.size = Pt(size); r.font.color.rgb = NAVY


def build(spec_path, out_path):
    spec = json.load(open(spec_path, encoding='utf-8'))
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    slides = spec['slides']
    meta = spec.get('meta', {})
    dividers = [x for x in slides if str(x.get('id', '')).startswith('section-')]
    div_index = {id(d): k + 1 for k, d in enumerate(dividers)}
    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        is_title = (i == 0) or s.get('id') == 'title'
        is_div = str(s.get('id', '')).startswith('section-')
        fig = resolve_fig(s.get('figure'))
        bullets = s.get('bullets') or []

        if is_title:
            ts = dict(s); ts.setdefault('author', meta.get('author'))
            build_title_slide(slide, ts)
            set_notes(slide, s.get('notes', ''))
            continue

        if is_div:
            build_divider_slide(slide, s, div_index[id(s)], len(dividers))
            set_notes(slide, s.get('notes', ''))
            continue

        if fig is None:
            build_textcard_slide(slide, s, bullets)
            set_notes(slide, s.get('notes', ''))
            continue

        add_title(slide, s['title'])

        if False:
            pass
        else:
            with Image.open(fig) as im:
                iw, ih = im.size
            aspect = iw / ih
            if aspect >= 1.5 and bullets:
                # wide figure: bullets band under title, image fills the rest
                add_bullets(slide, bullets, Inches(0.6), Inches(1.25),
                            Inches(12.1), Inches(1.55), size=16)
                l, t, w, h = fit_box(fig, Inches(0.5), Inches(2.95),
                                     Inches(12.33), Inches(4.3))
                slide.shapes.add_picture(str(fig), l, t, w, h)
            elif not bullets:
                # figure only: large centered
                l, t, w, h = fit_box(fig, Inches(0.5), Inches(1.3),
                                     Inches(12.33), Inches(5.9))
                slide.shapes.add_picture(str(fig), l, t, w, h)
            else:
                # tall/square: bullets left, image right
                add_bullets(slide, bullets, Inches(0.6), Inches(1.4),
                            Inches(4.5), Inches(5.6), size=18)
                l, t, w, h = fit_box(fig, Inches(5.3), Inches(1.3),
                                     Inches(7.6), Inches(5.9))
                slide.shapes.add_picture(str(fig), l, t, w, h)
            if s.get('video_placeholder') and fig is not None:
                add_play_button(slide, l, t, w, h)

        set_notes(slide, s.get('notes', ''))

    prs.save(out_path)
    print(f'  ✓ {out_path}  ({len(slides)} slides)')


if __name__ == '__main__':
    spec = sys.argv[1] if len(sys.argv) > 1 else str(HERE / 'deck_final.json')
    out = sys.argv[2] if len(sys.argv) > 2 else str(HERE / 'lab3_journey.pptx')
    build(spec, out)
